import sys
import os
import json
import re
import requests
import pandas as pd
import docx
import google.generativeai as genai
import streamlit as st
import matplotlib.pyplot as plt
import seaborn as sns
from pypdf import PdfReader
from youtube_transcript_api import YouTubeTranscriptApi
from bs4 import BeautifulSoup
from PIL import Image
from io import BytesIO
from datetime import date
from pptx import Presentation
from pptx.util import Inches, Pt
from docx.shared import Pt, RGBColor, Inches
from docx.enum.text import WD_ALIGN_PARAGRAPH

# --- 1. CONFIGURACIÓN E IDENTIDAD ---
st.set_page_config(
    page_title="IkigAI V1.86 - Executive Workstation", 
    page_icon="🧬", 
    layout="wide",
    initial_sidebar_state="expanded"
)

# Estilo CSS Zen: Contraste Quirúrgico y Ergonomía Móvil
st.markdown("""
    <style>
    @import url('https://fonts.googleapis.com/css2?family=Inter:wght@300;400;600&display=swap');
    .stApp { background-color: #000000 !important; font-family: 'Inter', sans-serif !important; }
    [data-testid="stSidebar"] { background-color: #080808 !important; border-right: 1px solid #1A1A1A !important; }
    [data-testid="stSidebar"] label, [data-testid="stSidebar"] p, [data-testid="stSidebar"] h1, h2, h3 { color: #FFFFFF !important; }
    [data-testid="stChatMessage"] { background-color: #050505 !important; border: 1px solid #1A1A1A !important; }
    .stMarkdown p, .stMarkdown li { color: #FFFFFF !important; font-size: 16px !important; line-height: 1.7 !important; }
    .stDownloadButton button, .stButton button { width: 100%; border-radius: 4px; background-color: transparent !important; color: #00E6FF !important; border: 1px solid #00E6FF !important; font-weight: 600; }
    .stDownloadButton button:hover, .stButton button:hover { background-color: #00E6FF !important; color: #000000 !important; }
    .section-tag { font-size: 11px; color: #666; letter-spacing: 1.5px; margin: 15px 0 5px 0; font-weight: 600; }
    .stExpander { border: 1px solid #1A1A1A !important; background-color: #050505 !important; border-radius: 8px !important; }
    textarea { background-color: #0D1117 !important; color: #FFFFFF !important; border: 1px solid #00E6FF !important; font-family: 'Courier New', monospace !important; font-size: 14px !important; }
    .stCheckbox { background-color: #111; padding: 5px; border-radius: 5px; border: 1px solid #333; margin-top: 10px; }
    </style>
""", unsafe_allow_html=True)

if "GOOGLE_API_KEY" in st.secrets:
    genai.configure(api_key=st.secrets["GOOGLE_API_KEY"])
else:
    st.error("🔑 Configure su API Key en st.secrets.")
    st.stop()

ROLES = {
    "Coach de Alto Desempeño": "ROI cognitivo y sostenibilidad administrativa.",
    "Director Centro Telemedicina": "Innovación, IA y Salud Digital UNAL.",
    "Vicedecano Académico": "Gestión académica, normativa y MD-PhD.",
    "Director de UCI": "Rigor clínico, datos HUN y seguridad.",
    "Investigador Científico": "Metodología, rigor y APA 7.",
    "Consultor Salud Digital": "BID/MinSalud y territorio.",
    "Professor Universitario": "Pedagogía médica disruptiva.",
    "Estratega de Trading": "Gestión de riesgo y SMC."
}

# --- 2. FUNCIONES DE LECTURA Y PERSISTENCIA (AJUSTE NATIVO) ---
def get_pdf_text(f): return "".join([p.extract_text() for p in PdfReader(f).pages])
def get_docx_text(f): return "\n".join([p.text for p in docx.Document(f).paragraphs])
def get_excel_text(f): return pd.read_excel(f).to_string()

# --- CAMBIO CLAVE: REEMPLAZO DE LANGCHAIN POR BUSCADOR NATIVO ---
DB_NATIVO = "memoria_nativa.json"
DATA_FOLDER = "biblioteca_master"

def actualizar_memoria_persistente():
    if not os.path.exists(DATA_FOLDER): os.makedirs(DATA_FOLDER)
    biblioteca_data = []
    archivos_encontrados = 0

    for file in os.listdir(DATA_FOLDER):
        if file.lower().endswith(".pdf"):
            try:
                with open(os.path.join(DATA_FOLDER, file), "rb") as f:
                    texto = get_pdf_text(f)
                    if texto and texto.strip():
                        # Segmentación Nativa: Chunks de 1200 caracteres (Reemplaza LangChain)
                        chunks = [texto[i:i+1200] for i in range(0, len(texto), 1000)]
                        for c in chunks:
                            biblioteca_data.append({"content": c, "source": file})
                        archivos_encontrados += 1
            except: continue

    if archivos_encontrados == 0:
        return "⚠️ Carpeta vacía o sin archivos legibles."

    with open(DB_NATIVO, "w", encoding="utf-8") as f:
        json.dump(biblioteca_data, f, ensure_ascii=False)
    
    return f"✅ ÉXITO: {archivos_encontrados} documentos sincronizados nativamente."

def exportar_sesion():
    mensajes_finales = []
    for i, msg in enumerate(st.session_state.messages):
        nuevo_msg = msg.copy()
        if msg["role"] == "assistant" and f"edit_{i}" in st.session_state:
            nuevo_msg["content"] = st.session_state[f"edit_{i}"]
        mensajes_finales.append(nuevo_msg)
    data = {"biblioteca": st.session_state.biblioteca, "messages": mensajes_finales, "last_analysis": st.session_state.last_analysis}
    return json.dumps(data, indent=4)

def cargar_sesion(json_data):
    data = json.loads(json_data)
    st.session_state.biblioteca = data["biblioteca"]
    st.session_state.messages = data["messages"]
    st.session_state.last_analysis = data["last_analysis"]

# --- 3. MOTOR DE EXPORTACIÓN (INTACTO) ---
def clean_markdown(text):
    text = re.sub(r'\*+', '', text)
    text = re.sub(r'^#+\s*', '', text, flags=re.MULTILINE)
    return text.strip()

def extraer_titulo_dictado(messages, indices_seleccionados):
    if not indices_seleccionados: return "MANUAL TÉCNICO DE TELESALUD"
    primer_contenido = messages[indices_seleccionados[0]]["content"]
    lineas = [l.strip() for l in primer_contenido.split('\n') if l.strip()]
    for linea in lineas:
        if any(x in linea.upper() for x in ["COMO IKIGAI", "PRESENTO", "DOCTOR", "HOLA", "ESTIMADO"]): continue
        titulo_limpio = re.sub(r'^#+\s*', '', linea)
        if len(titulo_limpio) > 5: return titulo_limpio.upper()
    return "DOCUMENTO ESTRATÉGICO DE GESTIÓN"

def download_word_compilado(indices_seleccionados, messages, role):
    doc = docx.Document()
    style = doc.styles['Normal']
    style.font.name = 'Arial'
    style.font.size = Pt(11)
    titulo_final = extraer_titulo_dictado(messages, indices_seleccionados)
    t = doc.add_heading(titulo_final, 0)
    t.alignment = WD_ALIGN_PARAGRAPH.CENTER
    doc.add_paragraph("").add_run()
    autor_p = doc.add_paragraph()
    run_a = autor_p.add_run("Jairo Antonio Pérez Cely")
    run_a.bold = True
    autor_p.alignment = WD_ALIGN_PARAGRAPH.CENTER
    for idx in sorted(indices_seleccionados):
        content = messages[idx]["content"]
        doc.add_paragraph(clean_markdown(content)).alignment = WD_ALIGN_PARAGRAPH.JUSTIFY
        doc.add_page_break()
    bio = BytesIO(); doc.save(bio); return bio.getvalue()

def download_pptx(content, role):
    prs = Presentation()
    slide = prs.slides.add_slide(prs.slide_layouts[0])
    slide.shapes.title.text = "IKIGAI Executive HUB"
    slide.placeholders[1].text = f"Autor: Jairo Antonio Pérez Cely\n{role}\n{date.today()}"
    bio = BytesIO(); prs.save(bio); return bio.getvalue()

def download_excel(content):
    try:
        lines = content.split('\n')
        table_data = []
        for line in lines:
            if '|' in line:
                cells = [c.strip() for c in line.split('|') if c.strip()]
                if cells: table_data.append(cells)
        if len(table_data) > 1:
            df = pd.DataFrame(table_data[1:], columns=table_data[0])
            bio = BytesIO()
            with pd.ExcelWriter(bio, engine='xlsxwriter') as writer:
                df.to_excel(writer, index=False, sheet_name='Datos_IkigAI')
            return bio.getvalue()
    except: return None

# --- 4. LÓGICA DE ESTADO (INTACTO) ---
if "biblioteca" not in st.session_state: st.session_state.biblioteca = {rol: "" for rol in ROLES.keys()}
if "messages" not in st.session_state: st.session_state.messages = []
if "last_analysis" not in st.session_state: st.session_state.last_analysis = ""
if "export_pool" not in st.session_state: st.session_state.export_pool = []
if "editor_version" not in st.session_state: st.session_state.editor_version = 0

# --- 5. BARRA LATERAL (INTACTO) ---
with st.sidebar:
    st.markdown("<h1 style='text-align: center; color: #00E6FF; font-size: 40px;'>🧬</h1>", unsafe_allow_html=True)
    st.markdown("<h2 style='text-align: center; letter-spacing: 5px; font-size: 24px;'>IKIGAI</h2>", unsafe_allow_html=True)
    st.divider()
    c1, c2 = st.columns(2)
    with c1:
        if st.button("🗑️ Reiniciar"):
            st.session_state.messages = []; st.session_state.export_pool = []; st.rerun()
    with c2:
        st.download_button("💾 Guardar", data=exportar_sesion(), file_name=f"Sesion_{date.today()}.json")
    st.divider()
    rol_activo = st.radio("Rol:", options=list(ROLES.keys()), label_visibility="collapsed")
    pool_actual = st.session_state.export_pool
    if pool_actual:
        st.divider()
        st.download_button("📄 Generar Word", data=download_word_compilado(pool_actual, st.session_state.messages, rol_activo), file_name=f"Reporte_{date.today()}.docx", use_container_width=True)
    
    st.divider()
    tab_doc, tab_url, tab_img = st.tabs(["📄 DOC/PPT", "🔗 URL", "🖼️ IMG"])
    with tab_doc:
        up = st.file_uploader("Subir archivos:", type=['pdf', 'docx', 'pptx'], accept_multiple_files=True)
        if st.button("🧠 Procesar archivos", use_container_width=True):
            raw_text = ""
            for f in up:
                if f.type == "application/pdf": raw_text += get_pdf_text(f)
            with st.spinner("Refinando..."):
                try:
                    refiner = genai.GenerativeModel('gemini-2.5-flash')
                    prompt_res = f"Extrae datos clave: {raw_text[:40000]}"
                    resumen = refiner.generate_content(prompt_res)
                    st.session_state.biblioteca[rol_activo] = resumen.text
                    st.success("Cargado.")
                except: st.session_state.biblioteca[rol_activo] = raw_text[:30000]

    if st.button("🧠 Sincronizar memoria máster", use_container_width=True):
        st.info(actualizar_memoria_persistente())

# --- 6. PANEL CENTRAL (AJUSTE RAG NATIVO) ---
ver = st.session_state.editor_version
for i, msg in enumerate(st.session_state.messages):
    with st.chat_message(msg["role"]):
        st.markdown(msg["content"])
        if msg["role"] == "assistant":
            if st.checkbox("📥 Incluir", key=f"sel_{i}_{ver}", value=(i in st.session_state.export_pool)):
                if i not in st.session_state.export_pool: st.session_state.export_pool.append(i); st.rerun()
            elif i in st.session_state.export_pool: st.session_state.export_pool.remove(i); st.rerun()

if pr := st.chat_input("Nuestro reto para hoy..."):
    from datetime import datetime
    fecha_actual = datetime.now().strftime("%d de %B de %Y")
    st.session_state.messages.append({"role": "user", "content": pr})
    with st.chat_message("user"): st.markdown(pr)
    
    with st.chat_message("assistant"):
        try:
            # --- CAMBIO CLAVE: BUSCADOR NATIVO ---
            contexto_rag = ""
            if os.path.exists(DB_NATIVO):
                with open(DB_NATIVO, "r", encoding="utf-8") as f:
                    data_master = json.load(f)
                palabras_clave = pr.lower().split()
                matches = []
                for item in data_master:
                    score = sum(1 for p in palabras_clave if p in item["content"].lower())
                    if score > 0: matches.append((score, item))
                matches.sort(key=lambda x: x[0], reverse=True)
                contexto_rag = "\n\n".join([f"(Fuente: {m[1]['source']}) {m[1]['content']}" for m in matches[:3]])

            contexto_reciente = st.session_state.biblioteca.get(rol_activo, "")
            model = genai.GenerativeModel('gemini-2.5-flash')
            mindset_seleccionado = ROLES.get(rol_activo, "Visión estratégica.")
            
            sys_prompt = f"""
            Actúa como {rol_activo}. Fecha: {fecha_actual}. Mindset: {mindset_seleccionado}
            CONOCIMIENTO RECIENTE (Sidebar): {contexto_reciente[:1000] if contexto_reciente else "N/A"}
            MEMORIA MÁSTER: {contexto_rag[:1000] if contexto_rag else "N/A"}

            ESTRUCTURA OBLIGATORIA:
            1. ### Triage Estratégico
            2. ### ROI Cognitivo
            3. ### Análisis multidimensional
            4. ### Propuesta táctica
            5. **Pregunta de Punto Ciego**
            """
            
            resp = model.generate_content([sys_prompt, pr])
            respuesta_final = resp.text
            if "Punto Ciego" not in respuesta_final:
                respuesta_final += f"\n\n---\n**Pregunta de Punto Ciego:** ¿Cómo afecta esto al ROI cognitivo?"

            st.markdown(respuesta_final)
            st.session_state.messages.append({"role": "assistant", "content": respuesta_final})
            st.rerun()
        except Exception as e:
            st.error(f"Error: {e}")
